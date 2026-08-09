from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, Response
import requests
import httpx
import traceback
import os
import base64
import html
import re
import time
from datetime import date
from urllib.parse import quote
from pathlib import Path
from dotenv import load_dotenv
from io import BytesIO
from docx import Document
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import BaseDocTemplate, Flowable, Frame, Image, KeepTogether, ListFlowable, ListItem, PageBreak, PageTemplate, Paragraph, Spacer, Table, TableStyle, XPreformatted
from reportlab.platypus.tableofcontents import TableOfContents
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor

load_dotenv()

from backend.github_oauth import router as github_router
from backend.ui_generator import router as ui_router
from backend.workspace import router as workspace_router
from backend.assistant import router as assistant_router

app = FastAPI(title='AOS API')
app.add_middleware(
    CORSMiddleware,
    allow_origins=['*'],
    allow_credentials=False,
    allow_methods=['*'],
    allow_headers=['*'],
)

app.include_router(github_router)
app.include_router(ui_router)
app.include_router(workspace_router)
app.include_router(assistant_router)



@app.exception_handler(HTTPException)
async def aos_http_exception_handler(request: Request, exception: HTTPException):
    """Retain the error response shape used by the existing browser scripts."""
    return JSONResponse(status_code=exception.status_code, content={'error': exception.detail})
APP_ROOT = Path(__file__).resolve().parent
PUBLIC_FILE_EXTENSIONS = {'.html', '.js', '.css', '.png', '.jpg', '.jpeg', '.svg', '.ico'}

API_KEYS = {
    'gemini': os.environ.get('GEMINI_API_KEY'),
    'groq': os.environ.get('GROQ_API_KEY'),
    'openrouter': os.environ.get('OPENROUTER_API_KEY'),
    'mistral': os.environ.get('MISTRAL_API_KEY'),
    'cohere': os.environ.get('COHERE_API_KEY')
}

# Keep image generation isolated from chat: configure this separately in
# Railway Variables. It must never be exposed in frontend code.
HF_TOKEN = os.environ.get('HF_TOKEN')
HF_IMAGE_MODELS = {'black-forest-labs/FLUX.1-schnell'}

@app.get('/')
def index():
    return FileResponse(APP_ROOT / 'index.html')

@app.post('/api/chat')
def chat(data: dict):
    if not data:
        raise HTTPException(status_code=400, detail='Invalid JSON data provided.')

    model = data.get('model')
    message = data.get('message')

    if not model or not message:
        raise HTTPException(status_code=400, detail='Model and message are required.')

    try:
        if model == 'gemini':
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-3.6-flash:generateContent?key={API_KEYS['gemini']}"
            payload = {'contents': [{'parts': [{'text': message}]}]}
            response = requests.post(url, json=payload, timeout=45)
            response.raise_for_status()
            data = response.json()
            reply = data.get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', 'No response from Gemini')
            return {'reply': reply}

        elif model == 'groq':
            url = 'https://api.groq.com/openai/v1/chat/completions'
            headers = {'Authorization': f"Bearer {API_KEYS['groq']}"}
            payload = {'model': 'llama-3.1-8b-instant', 'messages': [{'role': 'user', 'content': message}]}
            response = requests.post(url, headers=headers, json=payload, timeout=45)
            response.raise_for_status()
            data = response.json()
            reply = data.get('choices', [{}])[0].get('message', {}).get('content', 'No response from Groq')
            return {'reply': reply}

        elif model == 'openrouter':
            url = 'https://openrouter.ai/api/v1/chat/completions'
            headers = {'Authorization': f"Bearer {API_KEYS['openrouter']}"}
            payload = {'model': 'openai/gpt-3.5-turbo', 'messages': [{'role': 'user', 'content': message}]}
            response = requests.post(url, headers=headers, json=payload, timeout=45)
            response.raise_for_status()
            data = response.json()
            reply = data.get('choices', [{}])[0].get('message', {}).get('content', 'No response from Open Router')
            return {'reply': reply}

        elif model == 'mistral':
            url = 'https://api.mistral.ai/v1/chat/completions'
            headers = {'Authorization': f"Bearer {API_KEYS['mistral']}"}
            payload = {'model': 'mistral-small-latest', 'messages': [{'role': 'user', 'content': message}]}
            response = requests.post(url, headers=headers, json=payload, timeout=45)
            response.raise_for_status()
            data = response.json()
            reply = data.get('choices', [{}])[0].get('message', {}).get('content', 'No response from Mistral')
            return {'reply': reply}

        elif model == 'cohere':
            url = 'https://api.cohere.com/v1/chat'
            headers = {'Authorization': f"Bearer {API_KEYS['cohere']}"}
            payload = {'message': message, 'model': 'command-r-plus-08-2024'}
            response = requests.post(url, headers=headers, json=payload, timeout=45)
            response.raise_for_status()
            data = response.json()
            reply = data.get('text', 'No response from Cohere')
            return {'reply': reply}

        else:
            raise HTTPException(status_code=400, detail='Unsupported model')

    except requests.exceptions.RequestException as e:
        print("Request failed:", e)
        if hasattr(e, 'response') and e.response is not None:
            print("Response content:", e.response.content)
        raise HTTPException(status_code=502, detail='Failed to communicate with LLM provider API.')
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail='An internal server error occurred.')

@app.post('/api/generate/image')
async def generate_image(data: dict):
    """Generate an image through Hugging Face; the token never reaches the browser."""
    prompt = (data.get('prompt') or '').strip()
    model = data.get('model') or 'black-forest-labs/FLUX.1-schnell'
    layout_mode = data.get('layout_mode') or False
    
    if not prompt:
        raise HTTPException(status_code=400, detail='An image prompt is required.')
    if model not in HF_IMAGE_MODELS:
        raise HTTPException(status_code=400, detail='Unsupported Hugging Face image model.')
    if not HF_TOKEN:
        raise HTTPException(status_code=503, detail='Hugging Face image API is not configured on the server.')

    if layout_mode:
        gemini_api_key = os.getenv("GEMINI_API_KEY_2") or os.getenv("GEMINI_API_KEY")
        if not gemini_api_key:
            raise HTTPException(status_code=503, detail='Gemini API key is not configured on the server.')
        try:
            import sys
            backend_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend")
            if backend_dir not in sys.path:
                sys.path.append(backend_dir)
            from poster_generator import create_layout_aware_poster
            img_b64 = await create_layout_aware_poster(prompt, gemini_api_key, model)
            return {
                'image_url': img_b64,
                'provider': 'huggingface_layout_composite',
                'model': model
            }
        except Exception as error:
            import traceback
            traceback.print_exc()
            raise HTTPException(status_code=502, detail=f'Layout-aware image composition failed: {str(error)}')

    try:
        response = httpx.post(
            'https://router.huggingface.co/nscale/v1/images/generations',
            headers={'Authorization': f'Bearer {HF_TOKEN}'},
            json={'model': model, 'prompt': prompt, 'response_format': 'b64_json'},
            timeout=120,
        )
        response.raise_for_status()
        payload = response.json()
        result = (payload.get('data') or [{}])[0]
        encoded = result.get('b64_json')
        if not encoded:
            raise ValueError('Hugging Face returned no image data.')
        return {
            'image_url': f'data:image/png;base64,{encoded}',
            'provider': 'huggingface',
            'model': model
        }
    except Exception as hf_error:
        print(f'HF image failed, trying Pollinations fallback: {hf_error}')
        # Fallback to Pollinations (free, no auth required)
        try:
            from urllib.parse import quote
            import requests as sync_requests
            safe_prompt = quote(prompt)
            poll_url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width=1024&height=1024&nologo=true"
            poll_res = sync_requests.get(poll_url, timeout=90)
            if poll_res.status_code == 200 and len(poll_res.content) > 1000:
                encoded = base64.b64encode(poll_res.content).decode('utf-8')
                return {
                    'image_url': f'data:image/jpeg;base64,{encoded}',
                    'provider': 'pollinations',
                    'model': 'pollinations-ai'
                }
            raise ValueError(f'Pollinations returned status {poll_res.status_code}')
        except Exception as poll_error:
            print(f'Pollinations fallback also failed: {poll_error}')
            if isinstance(hf_error, httpx.HTTPError):
                detail = hf_error.response.text[:300] if hf_error.response is not None else str(hf_error)
                raise HTTPException(status_code=502, detail=f'Image generation failed (HF + Pollinations): {detail}')
            raise HTTPException(status_code=502, detail=f'Image generation failed: {str(hf_error)}')


DOCUMENT_TYPES = {'PDF', 'Word', 'Excel', 'PowerPoint', 'Resume', 'Invoice', 'Business Plan', 'Research'}


def document_plan(prompt, document_type):
    """Create a real document outline from the user's request (no mock content)."""
    title = ' '.join(prompt.split())[:96]
    return {
        'title': title,
        'summary': f'This {document_type.lower()} was prepared from the following request: {prompt}',
        'sections': [
            ('Purpose', f'The purpose of this {document_type.lower()} is to address: {prompt}'),
            ('Scope', 'Define the intended audience, deliverables, and boundaries before execution.'),
            ('Key requirements', 'Capture the essential actions, dependencies, quality expectations, and success criteria.'),
            ('Next steps', 'Review this draft, add project-specific facts, assign owners, and confirm the delivery timeline.')
        ]
    }


def safe_filename(value, fallback):
    cleaned = ''.join(char if char.isalnum() else '-' for char in value.lower()).strip('-')
    return (cleaned[:48] or fallback)


def create_docx(plan, document_type):
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    heading = document.add_heading(plan['title'], 0)
    heading.runs[0].font.color.rgb = None
    document.add_paragraph(document_type, style='Subtitle')
    document.add_paragraph(plan['summary'])
    for heading_text, body in plan['sections']:
        document.add_heading(heading_text, level=1)
        document.add_paragraph(body)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue(), 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'docx'


def create_xlsx(plan):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = 'Document Plan'
    sheet.append(['Section', 'Details'])
    sheet.append(['Title', plan['title']])
    sheet.append(['Summary', plan['summary']])
    for heading, body in plan['sections']:
        sheet.append([heading, body])
    for cell in sheet[1]:
        cell.font = Font(bold=True, color='FFFFFF')
        cell.fill = PatternFill('solid', fgColor='2563EB')
    sheet.column_dimensions['A'].width = 24
    sheet.column_dimensions['B'].width = 88
    for row in sheet.iter_rows(min_row=2):
        row[1].alignment = row[1].alignment.copy(wrap_text=True, vertical='top')
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue(), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'xlsx'


class AOSReportTemplate(BaseDocTemplate):
    """Platypus document template with a working table of contents."""
    def afterFlowable(self, flowable):
        if isinstance(flowable, Paragraph) and getattr(flowable, '_toc_level', None) is not None:
            key = f'chapter-{self.seq.nextf("chapter")}'
            self.canv.bookmarkPage(key)
            self.notify('TOCEntry', (flowable._toc_level, flowable.getPlainText(), self.page, key))


class RoundedNote(Flowable):
    def __init__(self, text, width):
        super().__init__()
        self.width = width
        self.paragraph = Paragraph(html.escape(text), ParagraphStyle('NoteBody', fontName='Helvetica', fontSize=10.5, leading=15, textColor=colors.HexColor('#1F2937')))
        self.height = 0

    def wrap(self, available_width, available_height):
        self.width = min(self.width, available_width)
        _, paragraph_height = self.paragraph.wrap(self.width - 28, available_height)
        self.height = paragraph_height + 26
        return self.width, self.height

    def draw(self):
        self.canv.setFillColor(colors.HexColor('#EFF6FF'))
        self.canv.setStrokeColor(colors.HexColor('#BFDBFE'))
        self.canv.roundRect(0, 0, self.width, self.height, 8, fill=1, stroke=1)
        self.paragraph.drawOn(self.canv, 14, 13)


def escape_text(value):
    return html.escape(str(value)).replace('\n', '<br/>')


def report_styles():
    base = getSampleStyleSheet()
    return {
        'cover_title': ParagraphStyle('CoverTitle', parent=base['Title'], fontName='Helvetica-Bold', fontSize=32, leading=39, alignment=TA_CENTER, textColor=colors.HexColor('#1F2937'), spaceAfter=16),
        'cover_subtitle': ParagraphStyle('CoverSubtitle', parent=base['Normal'], fontName='Helvetica', fontSize=13, leading=19, alignment=TA_CENTER, textColor=colors.HexColor('#64748B')),
        'heading1': ParagraphStyle('Heading1AOS', parent=base['Heading1'], fontName='Helvetica-Bold', fontSize=20, leading=26, textColor=colors.HexColor('#2563EB'), spaceBefore=0, spaceAfter=14, keepWithNext=True),
        'heading2': ParagraphStyle('Heading2AOS', parent=base['Heading2'], fontName='Helvetica-Bold', fontSize=16, leading=21, textColor=colors.HexColor('#1F2937'), spaceBefore=10, spaceAfter=8, keepWithNext=True),
        'heading3': ParagraphStyle('Heading3AOS', parent=base['Heading3'], fontName='Helvetica-Bold', fontSize=14, leading=18, textColor=colors.HexColor('#1F2937'), spaceBefore=8, spaceAfter=6, keepWithNext=True),
        'body': ParagraphStyle('BodyAOS', parent=base['BodyText'], fontName='Helvetica', fontSize=11.2, leading=17, textColor=colors.HexColor('#374151'), spaceAfter=11),
        'toc': ParagraphStyle('TOCAOS', parent=base['Normal'], fontName='Helvetica', fontSize=11, leading=16, textColor=colors.HexColor('#374151')),
        'code': ParagraphStyle('CodeAOS', fontName='Courier', fontSize=8.8, leading=12.5, textColor=colors.HexColor('#E5E7EB'), backColor=colors.HexColor('#111827'), borderPadding=11, borderRadius=6),
    }


def on_cover_page(canvas, document):
    canvas.saveState()
    canvas.setFillColor(colors.HexColor('#2563EB'))
    canvas.rect(0, A4[1] - 82, A4[0], 82, fill=1, stroke=0)
    canvas.setFillColor(colors.white)
    canvas.setFont('Helvetica-Bold', 12)
    canvas.drawString(48, A4[1] - 49, 'AOS  /  GENERATED DOCUMENT')
    canvas.setStrokeColor(colors.HexColor('#93C5FD'))
    canvas.setLineWidth(2)
    canvas.line(48, 48, A4[0] - 48, 48)
    canvas.setFillColor(colors.HexColor('#64748B'))
    canvas.setFont('Helvetica', 9)
    canvas.drawCentredString(A4[0] / 2, 31, f'Page {document.page}')
    canvas.restoreState()


def on_body_page(canvas, document):
    canvas.saveState()
    canvas.setStrokeColor(colors.HexColor('#2563EB'))
    canvas.setLineWidth(1.6)
    canvas.line(48, A4[1] - 44, A4[0] - 48, A4[1] - 44)
    canvas.setFillColor(colors.HexColor('#2563EB'))
    canvas.setFont('Helvetica-Bold', 9)
    canvas.drawString(48, A4[1] - 33, 'AOS DOCUMENT STUDIO')
    canvas.setFillColor(colors.HexColor('#64748B'))
    canvas.setFont('Helvetica', 8.5)
    canvas.drawRightString(A4[0] - 48, A4[1] - 33, date.today().strftime('%d %B %Y'))
    canvas.setStrokeColor(colors.HexColor('#CBD5E1'))
    canvas.setLineWidth(.7)
    canvas.line(48, 40, A4[0] - 48, 40)
    canvas.setFillColor(colors.HexColor('#64748B'))
    canvas.setFont('Helvetica', 8.5)
    canvas.drawString(48, 26, 'Prepared with AOS')
    canvas.drawRightString(A4[0] - 48, 26, f'Page {document.page}')
    canvas.restoreState()


def optional_images(image_payloads, max_width):
    flowables = []
    for data_url in image_payloads or []:
        if not isinstance(data_url, str) or not data_url.startswith('data:image/') or ',' not in data_url:
            continue
        try:
            image_bytes = base64.b64decode(data_url.split(',', 1)[1])
            image = Image(BytesIO(image_bytes))
            image._restrictSize(max_width, 4.6 * inch)
            flowables += [Spacer(1, 6), image, Spacer(1, 10)]
        except Exception:
            continue
    return flowables


def highlighted_code(code, style):
    highlighted_lines = []
    for line in code.splitlines():
        code_part, marker, comment = line.partition('#')
        escaped = html.escape(code_part)
        for keyword in ('def', 'class', 'return', 'import', 'from', 'if', 'else', 'for', 'while', 'True', 'False', 'None'):
            escaped = re.sub(rf'\b{keyword}\b', f'<font color="#60A5FA">{keyword}</font>', escaped)
        if marker:
            escaped += f'<font color="#94A3B8">#{html.escape(comment)}</font>'
        highlighted_lines.append(escaped)
    return XPreformatted('\n'.join(highlighted_lines), style)


def create_pdf(plan, document_type, options=None):
    options = options or {}
    buffer = BytesIO()
    page_width, page_height = A4
    margin = 52
    cover_frame = Frame(margin, 62, page_width - (margin * 2), page_height - 160, id='cover')
    body_frame = Frame(margin, 58, page_width - (margin * 2), page_height - 118, id='body')
    document = AOSReportTemplate(buffer, pagesize=A4, leftMargin=margin, rightMargin=margin, topMargin=58, bottomMargin=58)
    document.addPageTemplates([PageTemplate(id='Cover', frames=[cover_frame], onPage=on_cover_page), PageTemplate(id='Body', frames=[body_frame], onPage=on_body_page)])
    styles = report_styles()
    story = [Spacer(1, 2.1 * inch), Paragraph(escape_text(plan['title']), styles['cover_title']), Paragraph(f'{escape_text(document_type)}  •  Generated {date.today().strftime("%d %B %Y")}', styles['cover_subtitle']), Spacer(1, .45 * inch), RoundedNote('This professionally formatted document was generated from your request. Review and add verified project-specific details before formal submission.', page_width - (margin * 2)), PageBreak()]

    toc = TableOfContents()
    toc.levelStyles = [styles['toc'], ParagraphStyle('TOC2', parent=styles['toc'], leftIndent=18)]
    story += [Paragraph('Table of Contents', styles['heading1']), Spacer(1, 6), toc, PageBreak()]

    def chapter(title, body, level=0):
        heading = Paragraph(escape_text(title), styles['heading1'] if level == 0 else styles['heading2'])
        heading._toc_level = level
        return [heading, Paragraph(escape_text(body), styles['body'])]

    story += chapter('Executive Summary', plan['summary'])
    story += [RoundedNote('Important note: this document is a structured draft. Validate facts, dates, financial figures, and references before sharing externally.', page_width - (margin * 2)), Spacer(1, 14)]
    summary_table = Table([['Document type', document_type], ['Prepared', date.today().strftime('%d %B %Y')], ['Subject', plan['title']]], colWidths=[1.5 * inch, 4.6 * inch])
    summary_table.setStyle(TableStyle([('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#EFF6FF')), ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#1D4ED8')), ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'), ('FONTNAME', (1, 0), (1, -1), 'Helvetica'), ('GRID', (0, 0), (-1, -1), .5, colors.HexColor('#CBD5E1')), ('VALIGN', (0, 0), (-1, -1), 'TOP'), ('TOPPADDING', (0, 0), (-1, -1), 8), ('BOTTOMPADDING', (0, 0), (-1, -1), 8), ('LEFTPADDING', (0, 0), (-1, -1), 9)]))
    story += [summary_table, PageBreak()]

    for index, (heading, body) in enumerate(plan['sections'], start=1):
        story += chapter(f'{index}. {heading}', body)
        if heading == 'Key requirements':
            bullets = ['Define measurable outcomes and acceptance criteria.', 'Assign an owner and due date to each deliverable.', 'Review risks, dependencies, and approval steps.']
            story += [ListFlowable([ListItem(Paragraph(item, styles['body'])) for item in bullets], bulletType='bullet', leftIndent=20), Spacer(1, 8)]
        if heading == 'Next steps':
            steps = ['Review the draft with stakeholders.', 'Confirm final requirements and supporting evidence.', 'Publish the approved version and record ownership.']
            story += [ListFlowable([ListItem(Paragraph(item, styles['body'])) for item in steps], bulletType='1', leftIndent=22), Spacer(1, 8)]
        if index < len(plan['sections']):
            story.append(PageBreak())

    supplied_images = optional_images(options.get('images'), page_width - (margin * 2))
    if supplied_images:
        story += [PageBreak()] + chapter('Supporting Images', 'Images supplied with the request are included below.') + supplied_images

    code_examples = options.get('code_examples') or []
    if code_examples:
        story += [PageBreak()] + chapter('Code Examples', 'The following code examples were supplied with the request.')
        for code in code_examples[:4]:
            if isinstance(code, str) and code.strip():
                story += [highlighted_code(code[:5000], styles['code']), Spacer(1, 12)]

    references = [item for item in (options.get('references') or []) if isinstance(item, str) and item.strip()]
    story += [PageBreak()] + chapter('References', 'The following references were supplied with this request.' if references else 'No external references were supplied with this request.')
    if references:
        story.append(ListFlowable([ListItem(Paragraph(escape_text(reference), styles['body'])) for reference in references], bulletType='bullet', leftIndent=20))
    story += [PageBreak()] + chapter('Conclusion', f'This {document_type.lower()} provides a structured foundation for “{plan["title"]}”. Complete a final review before submission to ensure it meets the required academic, company, or project standards.')

    document.multiBuild(story)
    return buffer.getvalue(), 'application/pdf', 'pdf'


def create_pptx(plan, slide_count, theme):
    presentation = Presentation()
    palette = {
        'modern': (30, 64, 175), 'minimal': (248, 250, 252), 'bold': (15, 23, 42),
        'corporate': (30, 58, 138), 'creative': (109, 40, 217)
    }
    background = palette.get(theme, palette['modern'])
    foreground = (15, 23, 42) if theme == 'minimal' else (255, 255, 255)
    slide_data = [('Overview', plan['summary'])] + plan['sections']
    while len(slide_data) < slide_count:
        slide_data.append(('Project action', f'Expand the plan for: {plan["title"]}'))
    for index, (heading, body) in enumerate(slide_data[:slide_count]):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        fill = slide.background.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(*background)
        title_box = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.65), PptInches(12), PptInches(0.8))
        title = title_box.text_frame.paragraphs[0]
        title.text = plan['title'] if index == 0 else heading
        title.font.size = PptPt(29); title.font.bold = True; title.font.color.rgb = RGBColor(*foreground)
        body_box = slide.shapes.add_textbox(PptInches(0.85), PptInches(1.8), PptInches(11.4), PptInches(4.8))
        body_paragraph = body_box.text_frame.paragraphs[0]
        body_paragraph.text = body
        body_paragraph.font.size = PptPt(19); body_paragraph.font.color.rgb = RGBColor(*foreground)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue(), 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'pptx'


@app.post('/api/generate/document')
def generate_document(data: dict):
    prompt = (data.get('prompt') or '').strip()
    document_type = data.get('document_type') or 'PDF'
    if not prompt:
        raise HTTPException(status_code=400, detail='A document prompt is required.')
    if document_type not in DOCUMENT_TYPES:
        raise HTTPException(status_code=400, detail='Unsupported document type.')
    plan = document_plan(prompt, document_type)
    try:
        if document_type == 'Excel':
            content, mimetype, extension = create_xlsx(plan)
        elif document_type in {'PDF', 'Invoice'}:
            content, mimetype, extension = create_pdf(plan, document_type, data)
        elif document_type == 'PowerPoint':
            from backend.ppt_generator import create_premium_pptx
            content, mimetype, extension = create_premium_pptx(prompt, 6, 'modern')
        else:
            content, mimetype, extension = create_docx(plan, document_type)
        filename = f'{safe_filename(plan["title"], "document")}.{extension}'
        return Response(content=content, media_type=mimetype, headers={'Content-Disposition': f'attachment; filename="{filename}"'})
    except Exception as error:
        raise HTTPException(status_code=500, detail=f'Document generation failed: {str(error)}')


@app.post('/api/generate/presentation')
def generate_presentation(data: dict):
    prompt = (data.get('prompt') or '').strip()
    theme = data.get('theme') or 'modern'
    try:
        slide_count = int(data.get('slides') or 8)
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail='Slides must be a number.')
    if not prompt:
        raise HTTPException(status_code=400, detail='A presentation prompt is required.')
    if slide_count < 1 or slide_count > 20:
        raise HTTPException(status_code=400, detail='Choose between 1 and 20 slides.')
    plan = document_plan(prompt, 'Presentation')
    try:
        from backend.ppt_generator import create_premium_pptx
        content, mimetype, extension = create_premium_pptx(prompt, slide_count, theme)
        filename = f'{safe_filename(plan["title"], "presentation")}.{extension}'
        return Response(content=content, media_type=mimetype, headers={'Content-Disposition': f'attachment; filename="{filename}"'})
    except Exception as error:
        raise HTTPException(status_code=500, detail=f'Presentation generation failed: {str(error)}')


@app.post('/api/deploy/netlify')
def deploy_netlify(data: dict):
    import zipfile
    import secrets
    import time
    from backend.github_oauth import rest
    
    token = (data.get('token') or '').strip()
    name = (data.get('name') or '').strip()
    project_id = data.get('project_id')
    repo = data.get('repo')
    cmd = data.get('cmd')
    pub_dir = data.get('dir')

    if not token:
        raise HTTPException(status_code=400, detail='Netlify deployment credential is required.')
    if not name:
        raise HTTPException(status_code=400, detail='Project/Site name is required.')

    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }

    def netlify_slug(value: str) -> str:
        """Convert a repository or user label into a valid Netlify subdomain."""
        value = (value or '').strip().lower()
        value = value.removesuffix('.git').split('/')[-1]
        value = re.sub(r'[^a-z0-9-]+', '-', value)
        value = re.sub(r'-{2,}', '-', value).strip('-')
        return value[:54] or 'aos-project'

    def create_site_with_fallback(create_body: dict, repository_url: str | None = None):
        """Retry a taken subdomain using the GitHub repository name and a safe suffix."""
        requested_name = netlify_slug(create_body.get('name'))
        repository_name = netlify_slug(repository_url or '')
        candidates = [requested_name]
        if repository_name and repository_name not in candidates:
            candidates.append(repository_name)
        # A suffix guarantees the final retry is unique without overwriting another site.
        candidates.append(f"{repository_name or requested_name}-{secrets.token_hex(3)}")

        last_response = None
        for candidate in candidates:
            body = {**create_body, 'name': candidate}
            response = requests.post(
                'https://api.netlify.com/api/v1/sites',
                headers=headers,
                json=body,
                timeout=30
            )
            if response.ok:
                return response

            last_response = response
            if response.status_code in {401, 403}:
                raise HTTPException(
                    status_code=response.status_code,
                    detail='Netlify rejected this access token. Create a new Personal Access Token in Netlify User settings → Applications → Personal access tokens, then paste it into AOS.'
                )
            if response.status_code == 429:
                retry_after = response.headers.get('Retry-After')
                wait_message = f' Wait {retry_after} seconds before trying again.' if retry_after else ' Wait a few minutes before trying again.'
                raise HTTPException(
                    status_code=429,
                    detail='Netlify temporarily rate-limited site creation for this application.' + wait_message
                )
            error_text = response.text.lower()
            is_taken_subdomain = response.status_code == 422 and 'subdomain' in error_text and 'unique' in error_text
            if not is_taken_subdomain:
                break

        detail = last_response.text if last_response is not None else 'Netlify site creation failed.'
        raise HTTPException(
            status_code=last_response.status_code if last_response is not None else 502,
            detail=f"Netlify site creation failed after trying a GitHub repository fallback: {detail}"
        )

    def wait_for_netlify_deploy(site_id: str) -> dict:
        """Wait briefly for Netlify's actual build result; never report a pending site as ready."""
        deadline = time.monotonic() + 75
        latest = {}
        while time.monotonic() < deadline:
            response = requests.get(
                f'https://api.netlify.com/api/v1/sites/{site_id}/deploys?per_page=1',
                headers=headers,
                timeout=20
            )
            if response.ok:
                deployments = response.json() or []
                if deployments:
                    latest = deployments[0]
                    state = (latest.get('state') or '').lower()
                    if state == 'ready':
                        return latest
                    if state in {'error', 'failed'}:
                        summary = latest.get('error_message') or latest.get('summary') or 'Netlify build failed.'
                        clone_error = summary.lower()
                        if any(marker in clone_error for marker in ('unable to access repository', 'host key verification', 'could not read from remote repository')):
                            raise HTTPException(
                                status_code=502,
                                detail=(
                                    'Netlify cannot access this GitHub repository. In Netlify, open User settings → '
                                    'Applications → GitHub, connect GitHub, and grant Netlify access to this repository. '
                                    'Then deploy again.'
                                )
                            )
                        raise HTTPException(status_code=502, detail=f'Netlify build failed: {summary}')
            time.sleep(3)
        if not latest:
            raise HTTPException(
                status_code=502,
                detail=(
                    'Netlify created the site but did not start a build. Connect GitHub in '
                    'Netlify (User settings → Applications → GitHub), then deploy again.'
                )
            )
        # Netlify can need more than a minute for a first Git-based build.
        # Return its real pending state rather than a false "successful" result.
        return latest

    try:
        if project_id:
            files = rest("project_files", params={"project_id": f"eq.{project_id}", "select": "path,content"})
            if not files:
                raise HTTPException(status_code=404, detail='No project files found to deploy.')

            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'a', zipfile.ZIP_DEFLATED, False) as zip_file:
                for file in files:
                    path = file.get('path', '').lstrip('/')
                    content = file.get('content', '')
                    if path:
                        zip_file.writestr(path, content)
            
            zip_data = zip_buffer.getvalue()

            create_response = create_site_with_fallback({'name': name}, repo)
            
            site_data = create_response.json()
            site_id = site_data.get('id')

            deploy_headers = {
                'Authorization': f'Bearer {token}',
                'Content-Type': 'application/zip'
            }
            deploy_response = requests.post(
                f'https://api.netlify.com/api/v1/sites/{site_id}/deploys',
                headers=deploy_headers,
                data=zip_data,
                timeout=45
            )
            if not deploy_response.ok:
                raise HTTPException(status_code=deploy_response.status_code, detail=f"Netlify zip deployment failed: {deploy_response.text}")
            
            deploy_data = deploy_response.json()
            deployment = wait_for_netlify_deploy(site_id)
            site_url = site_data.get('ssl_url') or site_data.get('url') or deploy_data.get('ssl_url') or deploy_data.get('url')
            return {'ssl_url': site_url, 'status': (deployment.get('state') or 'building').lower()}
        else:
            clean_url = (repo or '').strip()
            clean_url = re.sub(r'^https?://', '', clean_url)
            clean_url = re.sub(r'^github\.com/*', '', clean_url)
            clean_url = re.sub(r'^com/*', '', clean_url)
            clean_url = re.sub(r'/+', '/', clean_url)
            clean_url = re.sub(r'\.git(hub)?$|\.gi$', '', clean_url)
            repo_path = clean_url.strip('/')
            if not re.fullmatch(r'[A-Za-z0-9_.-]+/[A-Za-z0-9_.-]+', repo_path):
                raise HTTPException(status_code=400, detail='Enter a valid public GitHub repository URL.')
            if cmd:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        'This repository needs a build command. Connect GitHub in Netlify to use Netlify builds, '
                        'or provide a repository that already contains its publish-ready files.'
                    )
                )

            # Netlify documents that its API repo object is complex and relies on a separate
            # Netlify GitHub App connection. For a public static repository, deploy the actual
            # repository files directly instead of asking Netlify to clone via that connection.
            # Do not call api.github.com here: Railway shares its outgoing IP and can hit
            # GitHub's anonymous API rate limit. Codeload is the public archive endpoint.
            archive_response = None
            for branch in ('main', 'master'):
                response = requests.get(
                    f'https://codeload.github.com/{repo_path}/zip/refs/heads/{branch}',
                    headers={'User-Agent': 'AOS-Netlify-Deploy'},
                    timeout=60
                )
                if response.ok:
                    archive_response = response
                    break
            if archive_response is None:
                raise HTTPException(
                    status_code=400,
                    detail='AOS could not download this public GitHub repository. Check the repository URL and ensure its default branch is main or master.'
                )

            # GitHub archives have a generated top-level folder. Remove it so index.html is
            # published at the Netlify site root. If a publish directory is supplied, use it.
            source_archive = BytesIO(archive_response.content)
            deploy_archive = BytesIO()
            selected_directory = (pub_dir or '').strip().strip('/')
            included_files = 0
            has_index_html = False
            with zipfile.ZipFile(source_archive) as source_zip, zipfile.ZipFile(deploy_archive, 'w', zipfile.ZIP_DEFLATED) as deploy_zip:
                for entry in source_zip.infolist():
                    if entry.is_dir():
                        continue
                    parts = entry.filename.split('/', 1)
                    if len(parts) != 2:
                        continue
                    relative_path = parts[1]
                    if selected_directory:
                        prefix = selected_directory + '/'
                        if not relative_path.startswith(prefix):
                            continue
                        relative_path = relative_path[len(prefix):]
                    if not relative_path:
                        continue
                    deploy_zip.writestr(relative_path, source_zip.read(entry))
                    included_files += 1
                    if relative_path.lower() == 'index.html':
                        has_index_html = True
            if not included_files:
                raise HTTPException(status_code=400, detail='No files were found in the selected Netlify publish directory.')
            if not has_index_html:
                location = selected_directory or 'the repository root'
                raise HTTPException(
                    status_code=400,
                    detail=(
                        f'No index.html was found in {location}. Direct Netlify deployment requires publish-ready '
                        'files. Build this project first and publish its dist folder, or add its complete package.json '
                        'and connect the repository through Netlify GitHub for a build deployment.'
                    )
                )

            create_response = create_site_with_fallback({'name': name}, repo)
            site_data = create_response.json()
            site_id = site_data.get('id')
            deploy_response = requests.post(
                f'https://api.netlify.com/api/v1/sites/{site_id}/deploys',
                headers={'Authorization': f'Bearer {token}', 'Content-Type': 'application/zip'},
                data=deploy_archive.getvalue(),
                timeout=90
            )
            if not deploy_response.ok:
                raise HTTPException(status_code=deploy_response.status_code, detail=f'Netlify file deployment failed: {deploy_response.text}')
            site_url = site_data.get('ssl_url') or site_data.get('url')
            deployment = wait_for_netlify_deploy(site_data.get('id'))
            return {'ssl_url': site_url, 'status': (deployment.get('state') or 'building').lower()}

    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post('/api/deploy/render')
def deploy_render(data: dict):
    token = (data.get('token') or '').strip()
    name = (data.get('name') or '').strip()
    repo = (data.get('repo') or '').strip()
    cmd = (data.get('cmd') or '').strip()
    pub_dir = (data.get('dir') or '').strip()
    service_type = (data.get('service_type') or 'static_site').strip().lower().replace(' ', '_')

    if not token:
        raise HTTPException(status_code=400, detail='Render API Key is required.')
    if not name:
        raise HTTPException(status_code=400, detail='Service Name is required.')
    if not repo:
        raise HTTPException(status_code=400, detail='GitHub Repository URL is required.')

    headers = {
        'Authorization': f'Bearer {token}',
        'Accept': 'application/json',
        'Content-Type': 'application/json'
    }

    # Step 1: Get owner ID from Render accounts
    try:
        owners_response = requests.get('https://api.render.com/v1/owners', headers=headers, timeout=20)
        if not owners_response.ok:
            if owners_response.status_code == 401:
                raise HTTPException(
                    status_code=401,
                    detail='Render rejected this API key. Create a new API key in Render Dashboard → Account Settings → API Keys, then paste that key into AOS.'
                )
            raise HTTPException(status_code=owners_response.status_code, detail=f"Render owner fetch failed: {owners_response.text}")
        owners = owners_response.json()
        if not owners or len(owners) == 0:
            raise HTTPException(status_code=400, detail="No owners found for this Render account.")
        owner_id = owners[0]['owner']['id']
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=f"Failed to fetch Render account owner ID: {str(e)}")

    # Render requires the type-specific configuration inside one `serviceDetails`
    # object.  The previous payload used `staticSiteDetails` at the top level,
    # which could create an incomplete service that never served the repository.
    service_kind = "web_service" if service_type in {"web_service", "node.js", "node"} else "static_site"
    body = {
        "type": service_kind,
        "name": name,
        "ownerId": owner_id,
        "repo": repo.removesuffix('.git'),
        "autoDeploy": "yes"
    }

    if service_kind == "static_site":
        # Do not force `npm run build`: plain HTML repositories do not have a
        # package.json.  Render receives a build command only when the user gave one.
        body["serviceDetails"] = {
            "publishPath": pub_dir or "."
        }
        if cmd:
            body["serviceDetails"]["buildCommand"] = cmd
    else:
        body["serviceDetails"] = {
            "runtime": "node",
            "envSpecificDetails": {
                "buildCommand": cmd or "npm install",
                "startCommand": "npm start"
            }
        }

    def service_url(service: dict) -> str:
        details = service.get('serviceDetails') or {}
        return details.get('url') or service.get('url') or ''

    def latest_deploy(service_id: str) -> dict:
        response = requests.get(
            f'https://api.render.com/v1/services/{service_id}/deploys?limit=1',
            headers=headers,
            timeout=20
        )
        if not response.ok:
            return {}
        result = response.json()
        if isinstance(result, list):
            item = result[0] if result else {}
        else:
            item = (result.get('items') or result.get('deploys') or [{}])[0]
        return item.get('deploy', item) if isinstance(item, dict) else {}

    def wait_for_render_deploy(service_id: str) -> dict:
        """Return Render's real outcome; never mark an in-progress build as live."""
        deadline = time.monotonic() + 70
        last = {}
        while time.monotonic() < deadline:
            last = latest_deploy(service_id)
            status = str(last.get('status') or last.get('state') or '').lower()
            if status in {'live', 'deployed', 'build_failed', 'failed', 'canceled', 'cancelled'}:
                return last
            time.sleep(5)
        return last

    try:
        create_response = requests.post('https://api.render.com/v1/services', headers=headers, json=body, timeout=30)
        duplicate_name = 'already in use' in create_response.text.lower()
        if create_response.status_code == 409 or duplicate_name:
            # A previous attempt can leave a service with the same name but an
            # invalid build configuration. Repair that service and trigger a
            # fresh deploy instead of forcing the user to delete it manually.
            services_response = requests.get(
                'https://api.render.com/v1/services',
                params={'ownerId': owner_id, 'limit': 100},
                headers=headers,
                timeout=20
            )
            services = services_response.json() if services_response.ok else []
            existing = next(
                (item.get('service', item) for item in services
                 if (item.get('service', item).get('name') == name)),
                None
            )
            if not existing or not existing.get('id'):
                raise HTTPException(status_code=409, detail='A Render service with this name already exists. Choose a different service name.')
            if existing.get('type') != service_kind:
                raise HTTPException(status_code=409, detail='A Render service with this name uses a different service type. Choose a different service name.')

            update_body = {
                'repo': body['repo'],
                'autoDeploy': body['autoDeploy'],
                'serviceDetails': body['serviceDetails']
            }
            update_response = requests.patch(
                f"https://api.render.com/v1/services/{existing['id']}",
                headers=headers,
                json=update_body,
                timeout=30
            )
            if not update_response.ok:
                raise HTTPException(status_code=update_response.status_code, detail=f"Render service update failed: {update_response.text}")
            service_data = update_response.json()
            trigger_response = requests.post(
                f"https://api.render.com/v1/services/{existing['id']}/deploys",
                headers=headers,
                json={'clearCache': 'do_not_clear'},
                timeout=30
            )
            if not trigger_response.ok:
                raise HTTPException(status_code=trigger_response.status_code, detail=f"Render deploy trigger failed: {trigger_response.text}")
        elif not create_response.ok:
            raise HTTPException(status_code=create_response.status_code, detail=f"Render service creation failed: {create_response.text}")
        else:
            created = create_response.json()
            # The create endpoint returns a service-and-deploy wrapper, while
            # update returns the service directly.
            service_data = created.get('service', created)
        service_id = service_data.get('id')
        if not service_id:
            raise HTTPException(status_code=502, detail='Render created no service ID. Please retry the deployment.')

        site_url = service_url(service_data)
        if not site_url:
            # Fetch the full service record instead of inventing a hostname.
            service_response = requests.get(f'https://api.render.com/v1/services/{service_id}', headers=headers, timeout=20)
            if service_response.ok:
                site_url = service_url(service_response.json())
        if not site_url:
            raise HTTPException(status_code=502, detail='Render did not return a public service URL.')

        deployment = wait_for_render_deploy(service_id)
        status = str(deployment.get('status') or deployment.get('state') or 'building').lower()
        if status in {'build_failed', 'failed', 'canceled', 'cancelled'}:
            raise HTTPException(
                status_code=502,
                detail='Render could not build this repository. Check the Render build logs and confirm the build command and publish path.'
            )
        return {'ssl_url': site_url, 'status': 'ready' if status in {'live', 'deployed'} else 'building', 'service_id': service_id}
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=f"Failed to create service on Render: {str(e)}")


@app.post('/api/deploy/render/status')
def render_deployment_status(data: dict):
    """Check an existing Render build without creating another deployment."""
    token = (data.get('token') or '').strip()
    service_id = (data.get('service_id') or '').strip()
    if not token or not service_id:
        raise HTTPException(status_code=400, detail='Render API Key and service ID are required.')

    headers = {'Authorization': f'Bearer {token}', 'Accept': 'application/json'}
    service_response = requests.get(f'https://api.render.com/v1/services/{service_id}', headers=headers, timeout=20)
    if service_response.status_code == 401:
        raise HTTPException(status_code=401, detail='Render rejected this API key.')
    if not service_response.ok:
        raise HTTPException(status_code=service_response.status_code, detail='Render service could not be found.')
    service = service_response.json()
    service_url = (service.get('serviceDetails') or {}).get('url') or service.get('url')

    deploy_response = requests.get(f'https://api.render.com/v1/services/{service_id}/deploys?limit=1', headers=headers, timeout=20)
    if not deploy_response.ok:
        raise HTTPException(status_code=deploy_response.status_code, detail='Render deployment status could not be loaded.')
    deploys = deploy_response.json()
    item = deploys[0] if isinstance(deploys, list) and deploys else {}
    deploy = item.get('deploy', item) if isinstance(item, dict) else {}
    raw_status = str(deploy.get('status') or deploy.get('state') or 'building').lower()
    if raw_status in {'live', 'deployed'}:
        status = 'ready'
    elif raw_status in {'build_failed', 'failed', 'canceled', 'cancelled'}:
        status = 'failed'
    else:
        status = 'building'
    return {'status': status, 'ssl_url': service_url, 'service_id': service_id}


@app.post('/api/deploy/railway')
def deploy_railway(data: dict):
    """Create a real Railway project/service from a GitHub repository."""
    token = (data.get('token') or '').strip()
    name = (data.get('name') or '').strip()
    repo = (data.get('repo') or '').strip()
    environment = (data.get('environment') or 'production').strip()
    if not token or not name or not repo:
        raise HTTPException(status_code=400, detail='Railway API token, project name, and GitHub repository URL are required.')

    repo_path = re.sub(r'^https?://github\.com/', '', repo, flags=re.I).removesuffix('.git').strip('/')
    if '/' not in repo_path:
        raise HTTPException(status_code=400, detail='Use a GitHub repository URL such as https://github.com/owner/repository.')

    headers = {'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'}

    def gql(query: str, variables: dict) -> dict:
        response = requests.post(
            'https://backboard.railway.com/graphql/v2',
            headers=headers,
            json={'query': query, 'variables': variables},
            timeout=30
        )
        if response.status_code == 401:
            raise HTTPException(status_code=401, detail='Railway rejected this API token. Create a new token in Railway Account Settings → Tokens.')
        if not response.ok:
            raise HTTPException(status_code=response.status_code, detail=f'Railway API request failed: {response.text}')
        result = response.json()
        if result.get('errors'):
            message = '; '.join(error.get('message', 'Railway request failed') for error in result['errors'])
            raise HTTPException(status_code=400, detail=f'Railway API: {message}')
        return result.get('data') or {}

    project_data = gql(
        'mutation projectCreate($input: ProjectCreateInput!) { projectCreate(input: $input) { id } }',
        {'input': {'name': name}}
    )
    project_id = (project_data.get('projectCreate') or {}).get('id')
    if not project_id:
        raise HTTPException(status_code=502, detail='Railway did not return a project ID.')

    project_info = gql(
        'query project($id: String!) { project(id: $id) { environments { edges { node { id name } } } } }',
        {'id': project_id}
    )
    environments = ((project_info.get('project') or {}).get('environments') or {}).get('edges') or []
    environment_id = next((edge['node']['id'] for edge in environments if edge.get('node', {}).get('name', '').lower() == environment.lower()), None)
    if not environment_id and environments:
        environment_id = environments[0].get('node', {}).get('id')
    if not environment_id:
        raise HTTPException(status_code=502, detail='Railway created the project but no deployable environment is available yet.')

    service_data = gql(
        'mutation serviceCreate($input: ServiceCreateInput!) { serviceCreate(input: $input) { id } }',
        {'input': {'projectId': project_id, 'name': name[:32], 'source': {'repo': repo_path}}}
    )
    service_id = (service_data.get('serviceCreate') or {}).get('id')
    if not service_id:
        raise HTTPException(status_code=502, detail='Railway did not return a service ID.')

    domain_data = gql(
        'mutation serviceDomainCreate($input: ServiceDomainCreateInput!) { serviceDomainCreate(input: $input) { domain } }',
        {'input': {'serviceId': service_id, 'environmentId': environment_id}}
    )
    domain = (domain_data.get('serviceDomainCreate') or {}).get('domain')
    if not domain:
        raise HTTPException(status_code=502, detail='Railway created the service but did not return a public domain.')

    # Railway automatically starts the first deployment after a GitHub-backed
    # service is created. Calling serviceInstanceDeploy immediately here can be
    # rejected before Railway has attached the GitHub source to the service.

    return {
        'ssl_url': f"https://{domain}" if not domain.startswith('http') else domain,
        'status': 'building',
        'project_id': project_id,
        'service_id': service_id
    }


@app.post('/api/deploy/railway/status')
def railway_deployment_status(data: dict):
    """Read Railway's latest status without creating another project or service."""
    token = (data.get('token') or '').strip()
    project_id = (data.get('project_id') or '').strip()
    service_id = (data.get('service_id') or '').strip()
    if not token or not project_id or not service_id:
        raise HTTPException(status_code=400, detail='Railway token, project ID, and service ID are required.')
    response = requests.post(
        'https://backboard.railway.com/graphql/v2',
        headers={'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'},
        json={
            'query': 'query deployments($input: DeploymentListInput!) { deployments(input: $input, first: 1) { edges { node { status } } } }',
            'variables': {'input': {'projectId': project_id, 'serviceId': service_id}}
        },
        timeout=30
    )
    if response.status_code == 401:
        raise HTTPException(status_code=401, detail='Railway rejected this API token.')
    if not response.ok:
        raise HTTPException(status_code=response.status_code, detail=f'Railway status request failed: {response.text}')
    result = response.json()
    if result.get('errors'):
        raise HTTPException(status_code=400, detail='Railway API: ' + '; '.join(error.get('message', 'status request failed') for error in result['errors']))
    edges = (((result.get('data') or {}).get('deployments') or {}).get('edges') or [])
    if not edges:
        return {'status': 'not_started', 'railway_status': 'NO_DEPLOYMENT', 'project_id': project_id, 'service_id': service_id}
    raw_status = str((edges[0].get('node') or {}).get('status') or 'BUILDING').upper()
    if raw_status == 'SUCCESS':
        status = 'ready'
    elif raw_status in {'FAILED', 'CRASHED', 'REMOVED', 'SKIPPED'}:
        status = 'failed'
    else:
        status = 'building'
    return {'status': status, 'railway_status': raw_status, 'project_id': project_id, 'service_id': service_id}


if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host='0.0.0.0', port=int(os.getenv('PORT', '8080')))
