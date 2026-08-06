"""
Premium PowerPoint Generator with AI-generated slide illustrations.
Uses LLM for rich slide content + Hugging Face for context-aware images.
"""
import os
import re
import json
import base64
import requests
import concurrent.futures
from io import BytesIO
from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import threading

HF_TOKEN = os.environ.get("HF_TOKEN")
GEMINI_IMAGE_KEY = os.environ.get("GEMINI_IMAGE_KEY") or os.environ.get("GEMINI_API_KEY_2") or os.environ.get("GEMINI_API_KEY")
_pollinations_lock = threading.Lock()
router = APIRouter()


class PPTRequest(BaseModel):
    prompt: str
    num_slides: int = 8
    theme: str = "professional"
    template: str = "business"

# ── Theme palettes ──────────────────────────────────────────────────────
THEMES = {
    'modern': {
        'bg': (15, 23, 42), 'fg': (255, 255, 255), 'accent': (59, 130, 246),
        'card': (30, 41, 59), 'muted': (148, 163, 184), 'gradient_end': (30, 64, 175)
    },
    'minimal': {
        'bg': (255, 255, 255), 'fg': (15, 23, 42), 'accent': (59, 130, 246),
        'card': (241, 245, 249), 'muted': (100, 116, 139), 'gradient_end': (226, 232, 240)
    },
    'bold': {
        'bg': (9, 9, 11), 'fg': (255, 255, 255), 'accent': (168, 85, 247),
        'card': (24, 24, 27), 'muted': (161, 161, 170), 'gradient_end': (88, 28, 135)
    },
    'corporate': {
        'bg': (15, 23, 42), 'fg': (255, 255, 255), 'accent': (14, 165, 233),
        'card': (30, 41, 59), 'muted': (148, 163, 184), 'gradient_end': (3, 105, 161)
    },
    'creative': {
        'bg': (30, 10, 60), 'fg': (255, 255, 255), 'accent': (244, 63, 94),
        'card': (50, 20, 80), 'muted': (196, 181, 253), 'gradient_end': (109, 40, 217)
    },
}


def _rgb(t):
    return RGBColor(*t)


def _generate_slide_content(prompt, slide_count):
    """Use LLM to generate rich, structured slide content as JSON.
    Attempts multiple Gemini models with backoff retry, and falls back to Groq/OpenRouter."""
    system = (
        "You are a professional presentation designer. Generate structured slide content as JSON.\n"
        "Return ONLY valid JSON (no markdown, no ```json wrapper). Structure:\n"
        "{\n"
        '  "title": "Presentation Title",\n'
        '  "slides": [\n'
        '    {\n'
        '      "title": "Slide Title",\n'
        '      "subtitle": "Optional subtitle",\n'
        '      "bullet_points": ["Point 1", "Point 2", "Point 3"],\n'
        '      "image_prompt": "Detailed description of a relevant illustration for this slide",\n'
        '      "speaker_notes": "Brief notes for the presenter"\n'
        '    }\n'
        "  ]\n"
        "}\n"
        f"Generate exactly {slide_count} slides. Make content detailed, professional, and directly relevant to the topic.\n"
        "The first slide should be the title/cover slide. The last slide should be a summary or Q&A slide.\n"
        "For image_prompt, describe a specific, high-quality illustration that matches the slide content."
    )
    user = f"Create a {slide_count}-slide presentation about: {prompt}"

    # Try Gemini first
    gemini_key = os.getenv("GEMINI_API_KEY_2") or os.getenv("GEMINI_API_KEY")
    if gemini_key:
        models = [
            os.getenv("PPT_GEMINI_MODEL", "gemini-3.6-flash"),
            "gemini-2.5-flash",
            "gemini-1.5-flash",
            "gemini-1.5-flash-8b"
        ]
        import time
        for model in models:
            for attempt in range(2):
                try:
                    response = requests.post(
                        f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent",
                        params={"key": gemini_key},
                        json={
                            "system_instruction": {"parts": [{"text": system}]},
                            "contents": [{"role": "user", "parts": [{"text": user}]}],
                            "generationConfig": {
                                "maxOutputTokens": 8192,
                                "responseMimeType": "application/json",
                            },
                        },
                        timeout=90,
                    )
                    if response.status_code == 200:
                        reply = response.json()["candidates"][0]["content"]["parts"][0]["text"]
                        match = re.search(r'(\{.*\})', reply.strip(), re.DOTALL)
                        if match:
                            generated = json.loads(match.group(1))
                            if isinstance(generated.get("slides"), list):
                                return generated
                    elif response.status_code in [429, 503]:
                        # Backoff retry
                        time.sleep(2 * (attempt + 1))
                        continue
                except Exception as e:
                    print(f"Gemini model {model} attempt {attempt+1} failed: {e}")
                    time.sleep(1)

    # Fallback to Groq
    groq_key = os.getenv("GROQ_API_KEY")
    if groq_key:
        print("Falling back to Groq for PPT content generation...")
        try:
            res = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers={"Authorization": f"Bearer {groq_key}", "Content-Type": "application/json"},
                json={
                    "model": "llama-3.3-70b-versatile",
                    "messages": [
                        {"role": "system", "content": system},
                        {"role": "user", "content": user}
                    ],
                    "response_format": {"type": "json_object"},
                    "temperature": 0.2
                },
                timeout=60
            )
            if res.status_code == 200:
                reply = res.json()["choices"][0]["message"]["content"]
                match = re.search(r'(\{.*\})', reply.strip(), re.DOTALL)
                if match:
                    generated = json.loads(match.group(1))
                    if isinstance(generated.get("slides"), list):
                        return generated
        except Exception as e:
            print(f"Groq fallback failed: {e}")

    # Fallback to OpenRouter
    or_key = os.getenv("OPENROUTER_API_KEY")
    if or_key:
        print("Falling back to OpenRouter for PPT content generation...")
        try:
            res = requests.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {or_key}",
                    "Content-Type": "application/json",
                    "HTTP-Referer": "https://aos-operating-system.web.app",
                    "X-Title": "AOS Studio"
                },
                json={
                    "model": "google/gemini-2.5-flash",
                    "messages": [
                        {"role": "system", "content": system},
                        {"role": "user", "content": user}
                    ],
                    "response_format": {"type": "json_object"}
                },
                timeout=60
            )
            if res.status_code == 200:
                reply = res.json()["choices"][0]["message"]["content"]
                match = re.search(r'(\{.*\})', reply.strip(), re.DOTALL)
                if match:
                    generated = json.loads(match.group(1))
                    if isinstance(generated.get("slides"), list):
                        return generated
        except Exception as e:
            print(f"OpenRouter fallback failed: {e}")

    raise RuntimeError("All LLM providers (Gemini, Groq, OpenRouter) failed to generate slide content.")


def _generate_image_gemini(prompt, width=400, height=300):
    """Generate image via Google Gemini Image API. Supports both generateContent (Nano Banana) and predict (Imagen)."""
    if not GEMINI_IMAGE_KEY:
        print("Gemini image generation skipped: GEMINI_IMAGE_KEY is not set.")
        return None

    # 1. Try generateContent with Nano Banana models (e.g. gemini-3.1-flash-image)
    nano_models = ["gemini-3.1-flash-image", "gemini-2.5-flash-image", "gemini-3-pro-image"]
    for model in nano_models:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={GEMINI_IMAGE_KEY}"
            payload = {
                "contents": [
                    {
                        "parts": [
                            {
                                "text": f"{prompt}, clean professional illustration, high quality, no text"
                            }
                        ]
                    }
                ]
            }
            res = requests.post(url, headers={"Content-Type": "application/json"}, json=payload, timeout=60)
            if res.status_code == 200:
                res_json = res.json()
                candidates = res_json.get("candidates", [])
                if candidates:
                    parts = candidates[0].get("content", {}).get("parts", [])
                    for part in parts:
                        inline_data = part.get("inlineData") or part.get("inline_data")
                        if inline_data:
                            b64 = inline_data.get("data")
                            if b64:
                                img_data = base64.b64decode(b64)
                                img = Image.open(BytesIO(img_data))
                                img = img.resize((width, height), Image.LANCZOS)
                                buf = BytesIO()
                                img.save(buf, format='PNG')
                                buf.seek(0)
                                return buf
            elif res.status_code == 429:
                print(f"Gemini image ({model}) rate limited or quota exceeded.")
        except Exception as e:
            print(f"Gemini image error for {model}: {type(e).__name__}: {e}")

    # 2. Try predict with Imagen models
    imagen_models = ["imagen-3.0-generate-002"]
    ratio = "1:1"
    if width > height * 1.3:
        ratio = "16:9"
    elif width > height * 1.1:
        ratio = "4:3"
    elif height > width * 1.3:
        ratio = "9:16"
    elif height > width * 1.1:
        ratio = "3:4"

    for model in imagen_models:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:predict?key={GEMINI_IMAGE_KEY}"
            payload = {
                "instances": [
                    {
                        "prompt": f"{prompt}, clean professional illustration, high quality, no text"
                    }
                ],
                "parameters": {
                    "sampleCount": 1,
                    "outputMimeType": "image/png",
                    "aspectRatio": ratio
                }
            }
            res = requests.post(url, headers={"Content-Type": "application/json"}, json=payload, timeout=60)
            if res.status_code == 200:
                res_json = res.json()
                predictions = res_json.get("predictions", [])
                if predictions:
                    b64 = predictions[0].get("bytesBase64Encoded")
                    if b64:
                        img_data = base64.b64decode(b64)
                        img = Image.open(BytesIO(img_data))
                        img = img.resize((width, height), Image.LANCZOS)
                        buf = BytesIO()
                        img.save(buf, format='PNG')
                        buf.seek(0)
                        return buf
        except Exception as e:
            print(f"Gemini image error for {model}: {type(e).__name__}: {e}")

    return None



def _generate_image_hf(prompt, width=400, height=300):
    """Try Hugging Face inference API. Returns BytesIO or None."""
    # Read at request time so the value always matches the Railway HF_TOKEN
    # used by AOS Image Generation after a service restart/redeploy.
    token = os.getenv("HF_TOKEN")
    if not token:
        print("PPT Hugging Face image skipped: HF_TOKEN is missing.")
        return None
    try:
        res = requests.post(
            'https://router.huggingface.co/nscale/v1/images/generations',
            headers={'Authorization': f'Bearer {token}'},
            json={
                'model': os.getenv('HF_IMAGE_MODEL', 'black-forest-labs/FLUX.1-schnell'),
                'prompt': f"{prompt}, clean professional illustration, high quality, no text",
                'response_format': 'b64_json'
            },
            timeout=90,
        )
        if res.status_code == 200:
            b64 = res.json().get('data', [{}])[0].get('b64_json')
            if b64:
                img_data = base64.b64decode(b64)
                img = Image.open(BytesIO(img_data))
                img = img.resize((width, height), Image.LANCZOS)
                buf = BytesIO()
                img.save(buf, format='PNG')
                buf.seek(0)
                return buf
        print(f"HF image failed: status {res.status_code} – {res.text[:200]}")
    except Exception as e:
        print(f"HF image error: {type(e).__name__}: {e}")
    return None


STABILITY_API_KEY = os.environ.get("STABILITY_API_KEY")

def _generate_image_stability(prompt, width=400, height=300):
    """Generate image via Stability AI Core API. Returns BytesIO or None."""
    if not STABILITY_API_KEY:
        print("Stability image skipped: STABILITY_API_KEY is not set.")
        return None
    try:
        ratio = "1:1"
        if width > height * 1.3:
            ratio = "16:9"
        elif height > width * 1.3:
            ratio = "9:16"
            
        res = requests.post(
            'https://api.stability.ai/v2beta/stable-image/generate/core',
            headers={'Authorization': f'Bearer {STABILITY_API_KEY}', 'accept': 'image/*'},
            files={'prompt': (None, f"{prompt}, clean professional illustration, high quality, no text")},
            data={'output_format': 'png', 'aspect_ratio': ratio},
            timeout=60
        )
        if res.status_code == 200:
            img = Image.open(BytesIO(res.content))
            img = img.resize((width, height), Image.LANCZOS)
            buf = BytesIO()
            img.save(buf, format='PNG')
            buf.seek(0)
            return buf
        print(f"Stability image failed: status {res.status_code} – {res.text[:200]}")
    except Exception as e:
        print(f"Stability image error: {type(e).__name__}: {e}")
    return None


def _generate_placeholder_gradient(theme_name, width=400, height=300):
    """Generate a premium gradient placeholder image matching the theme palette as a fallback."""
    try:
        from PIL import ImageDraw
        palette = THEMES.get(theme_name, THEMES['modern'])
        c1 = palette.get('accent', (59, 130, 246))
        c2 = palette.get('gradient_end', (30, 64, 175))
        
        # Create a linear gradient image
        img = Image.new("RGB", (width, height), c1)
        draw = ImageDraw.Draw(img)
        
        # Draw vertical gradient
        for y in range(height):
            # Interpolate between c1 and c2
            t = y / (height - 1)
            r = int(c1[0] + (c2[0] - c1[0]) * t)
            g = int(c1[1] + (c2[1] - c1[1]) * t)
            b = int(c1[2] + (c2[2] - c1[2]) * t)
            draw.line([(0, y), (width, y)], fill=(r, g, b))
            
        buf = BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"Error generating placeholder gradient: {e}")
        # Return a simple solid color image if gradient fails
        try:
            img = Image.new("RGB", (width, height), (30, 41, 59))
            buf = BytesIO()
            img.save(buf, format='PNG')
            buf.seek(0)
            return buf
        except Exception:
            return None


def _generate_image_pollinations(prompt, width=400, height=300):
    """Generate image via Pollinations AI. Returns BytesIO or None."""
    import time
    from urllib.parse import quote

    safe_prompt = quote(prompt)
    poll_url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width={width}&height={height}&nologo=true"
    
    with _pollinations_lock:
        # Sleep a little bit to space out requests to Pollinations and prevent rate limiting
        time.sleep(1.5)
        for attempt in range(3):
            try:
                res = requests.get(poll_url, timeout=60)
                if res.status_code == 200 and len(res.content) > 1000:
                    img = Image.open(BytesIO(res.content))
                    img = img.resize((width, height), Image.LANCZOS)
                    buf = BytesIO()
                    img.save(buf, format='PNG')
                    buf.seek(0)
                    return buf
                elif res.status_code == 429:
                    sleep_time = 4 + attempt * 3
                    print(f"Pollinations image rate limited (429), retrying in {sleep_time}s... (attempt {attempt + 1}/3)")
                    time.sleep(sleep_time)
                else:
                    print(f"Pollinations image failed: status {res.status_code}")
                    # Try once more after a small break
                    time.sleep(2)
            except Exception as e:
                print(f"Pollinations image error (attempt {attempt + 1}/3): {type(e).__name__}: {e}")
                time.sleep(2)
    return None


def _generate_image(prompt, width=400, height=300, theme='modern'):
    """Generate every PPT illustration through the AOS Hugging Face setup."""
    provider = "huggingface"
    buf = _generate_image_hf(prompt, width, height)
    if buf:
        return buf
    print(f"PPT image provider '{provider}' did not return an image; using a local visual fallback.")
    return _generate_placeholder_gradient(theme, width, height)


def _add_rounded_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    """Add a rounded rectangle shape as a card background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    shape.line.fill.background()  # no border
    shape.shadow.inherit = False
    # Bring to back so text overlays it
    return shape


def _add_text(slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(color)
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def create_premium_pptx_data(prompt, slide_count=8, theme='modern'):
    """Generate structured presentation slides and in-memory images, then return the PPTX bytes, slides text data, and slide images dict."""
    palette = THEMES.get(theme, THEMES['modern'])

    # 1. Generate structured slide content via LLM
    content = _generate_slide_content(prompt, slide_count)
    slides_data = content.get('slides', [])
    pres_title = content.get('title', prompt[:60])

    # Pad or trim to requested count
    while len(slides_data) < slide_count:
        slides_data.append({"title": "Additional Points", "subtitle": "", "bullet_points": [f"More details about {prompt}"], "image_prompt": f"illustration about {prompt}", "speaker_notes": ""})
    slides_data = slides_data[:slide_count]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # ── Pre-generate all images concurrently ──
    # Every slide receives an illustration by default. Set a smaller value only
    # when deliberately requesting a faster text-and-design draft.
    remote_image_count = max(0, min(int(os.getenv("PPT_REMOTE_IMAGE_COUNT", str(slide_count))), slide_count))
    remote_image_indexes = []
    if remote_image_count:
        remote_image_indexes.append(0)
    if remote_image_count > 1 and slide_count > 1:
        remote_image_indexes.append(slide_count - 1)
    if remote_image_count > 2:
        step = max(1, (slide_count - 2) // (remote_image_count - 2))
        remote_image_indexes.extend(range(1, slide_count - 1, step))
    remote_image_indexes = set(remote_image_indexes[:remote_image_count])

    image_prompts = []
    for idx, sd in enumerate(slides_data):
        if idx not in remote_image_indexes:
            continue
        is_cover = idx == 0
        is_last = idx == slide_count - 1
        if is_cover:
            image_prompts.append((idx, sd.get('image_prompt', prompt), 500, 350))
        elif is_last:
            image_prompts.append((idx, sd.get('image_prompt', f'thank you summary illustration about {prompt}'), 480, 400))
        else:
            image_prompts.append((idx, sd.get('image_prompt', f'illustration about {prompt}'), 480, 400))

    slide_images = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=min(4, len(image_prompts) or 1)) as executor:
        future_to_idx = {}
        for idx, img_prompt, w, h in image_prompts:
            if img_prompt:
                future = executor.submit(_generate_image, img_prompt, w, h, theme)
                future_to_idx[future] = idx
        for future in concurrent.futures.as_completed(future_to_idx):
            slide_idx = future_to_idx[future]
            try:
                slide_images[slide_idx] = future.result()
            except Exception as e:
                print(f"PPT image for slide {slide_idx} raised: {e}")
                slide_images[slide_idx] = None

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # ── Background ──
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(palette['bg'])

        # ── Accent bar (thin strip at top) ──
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(palette['accent'])
        bar.line.fill.background()

        is_cover = idx == 0
        is_last = idx == slide_count - 1

        if is_cover:
            # ══════════ COVER SLIDE ══════════
            _add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
                      pres_title, 44, palette['fg'], bold=True, alignment=PP_ALIGN.CENTER)

            subtitle = sd.get('subtitle', '')
            if subtitle:
                _add_text(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.8),
                          subtitle, 22, palette['muted'], alignment=PP_ALIGN.CENTER)

            divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.06))
            divider.fill.solid()
            divider.fill.fore_color.rgb = _rgb(palette['accent'])
            divider.line.fill.background()

            img_buf = slide_images.get(idx)
            if img_buf and idx in {0, len(slides_data) - 1}:
                slide.shapes.add_picture(img_buf, Inches(4.4), Inches(4.2), Inches(4.5), Inches(3.0))

        elif is_last:
            # ══════════ CLOSING SLIDE ══════════
            _add_text(slide, Inches(1.3), Inches(0.3), Inches(7), Inches(0.7),
                      sd.get('title', 'Thank You'), 28, palette['fg'], bold=True)

            subtitle = sd.get('subtitle', '')
            if subtitle:
                _add_text(slide, Inches(1.3), Inches(1.0), Inches(7), Inches(0.5),
                          subtitle, 16, palette['muted'])

            card_top = Inches(1.6)
            card_height = Inches(5.2)
            _add_rounded_rect(slide, Inches(0.5), card_top, Inches(7.5), card_height, palette['card'])

            bullets = sd.get('bullet_points', [])
            for bi, bp in enumerate(bullets[:6]):
                y_pos = card_top + Inches(0.4 + bi * 0.7)
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_pos + Inches(0.12), Inches(0.18), Inches(0.18))
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(palette['accent'])
                dot.line.fill.background()
                _add_text(slide, Inches(1.4), y_pos, Inches(6.2), Inches(0.6),
                          bp, 16, palette['fg'])

            img_buf = slide_images.get(idx)
            img_prompt = sd.get('image_prompt', f'thank you summary illustration about {prompt}')
            if img_buf:
                _add_rounded_rect(slide, Inches(8.3), card_top, Inches(4.5), card_height, palette['card'])
                slide.shapes.add_picture(img_buf, Inches(8.5), Inches(1.9), Inches(4.1), Inches(3.4))
                _add_text(slide, Inches(8.5), Inches(5.5), Inches(4.1), Inches(0.5),
                          img_prompt[:50], 11, palette['muted'], alignment=PP_ALIGN.CENTER)

        else:
            # ══════════ CONTENT SLIDES ══════════
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.35), Inches(0.55), Inches(0.55))
            badge.fill.solid()
            badge.fill.fore_color.rgb = _rgb(palette['accent'])
            badge.line.fill.background()
            badge_tf = badge.text_frame
            badge_tf.paragraphs[0].text = str(idx + 1)
            badge_tf.paragraphs[0].font.size = Pt(16)
            badge_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            badge_tf.paragraphs[0].font.bold = True
            badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            _add_text(slide, Inches(1.3), Inches(0.3), Inches(7), Inches(0.7),
                      sd.get('title', ''), 28, palette['fg'], bold=True)

            subtitle = sd.get('subtitle', '')
            if subtitle:
                _add_text(slide, Inches(1.3), Inches(1.0), Inches(7), Inches(0.5),
                          subtitle, 16, palette['muted'])

            card_top = Inches(1.6)
            card_height = Inches(5.2)
            _add_rounded_rect(slide, Inches(0.5), card_top, Inches(7.5), card_height, palette['card'])

            bullets = sd.get('bullet_points', [])
            for bi, bp in enumerate(bullets[:6]):
                y_pos = card_top + Inches(0.4 + bi * 0.7)
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_pos + Inches(0.12), Inches(0.18), Inches(0.18))
                dot.fill.solid()
                dot.fill.fore_color.rgb = _rgb(palette['accent'])
                dot.line.fill.background()
                _add_text(slide, Inches(1.4), y_pos, Inches(6.2), Inches(0.6),
                          bp, 16, palette['fg'])

            img_buf = slide_images.get(idx)
            img_prompt = sd.get('image_prompt', f'illustration about {prompt}')
            if img_buf:
                _add_rounded_rect(slide, Inches(8.3), card_top, Inches(4.5), card_height, palette['card'])
                slide.shapes.add_picture(img_buf, Inches(8.5), Inches(1.9), Inches(4.1), Inches(3.4))
                _add_text(slide, Inches(8.5), Inches(5.5), Inches(4.1), Inches(0.5),
                          img_prompt[:50], 11, palette['muted'], alignment=PP_ALIGN.CENTER)

        # ── Footer ──
        _add_text(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.4),
                  pres_title, 10, palette['muted'])
        _add_text(slide, Inches(10), Inches(7.0), Inches(3), Inches(0.4),
                  f"Slide {idx + 1} of {slide_count}", 10, palette['muted'], alignment=PP_ALIGN.RIGHT)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue(), slides_data, slide_images


def create_premium_pptx(prompt, slide_count=8, theme='modern'):
    """Generate a premium PowerPoint presentation bytes (backward compatible)."""
    pptx_bytes, slides_data, slide_images = create_premium_pptx_data(prompt, slide_count, theme)
    return pptx_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'pptx'


@router.post("/api/generate-ppt")
async def generate_ppt(request: PPTRequest):
    prompt = request.prompt.strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="Please enter a presentation topic.")

    slide_count = max(3, min(request.num_slides, 20))
    theme_aliases = {"professional": "corporate", "business": "corporate"}
    theme = theme_aliases.get(request.theme.lower(), request.theme.lower())
    if theme not in THEMES:
        theme = "corporate"

    try:
        pptx_bytes, slides_data, slide_images = create_premium_pptx_data(prompt, slide_count, theme)
        filename = re.sub(r"[^a-zA-Z0-9_-]+", "-", prompt[:60]).strip("-") or "aos-presentation"
        
        # Base64-encode the presentation bytes
        pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
        
        # Base64-encode the slide illustrations
        slides_response = []
        for idx, sd in enumerate(slides_data):
            img_buf = slide_images.get(idx)
            img_b64 = None
            if img_buf:
                # Return thumbnails for every preview slide. Full-size images
                # stay in the editable PPTX, keeping the HTTP response small.
                try:
                    image = Image.open(BytesIO(img_buf.getvalue())).convert("RGB")
                    image.thumbnail((240, 160), Image.LANCZOS)
                    thumbnail = BytesIO()
                    image.save(thumbnail, format="JPEG", quality=65, optimize=True)
                    img_b64 = base64.b64encode(thumbnail.getvalue()).decode('utf-8')
                except Exception:
                    img_b64 = None
                
            slides_response.append({
                "title": sd.get('title', ''),
                "subtitle": sd.get('subtitle', ''),
                "bullet_points": sd.get('bullet_points', []),
                "image_b64": img_b64
            })
            
        return {
            "filename": f"{filename}.pptx",
            "pptx": pptx_b64,
            "slides": slides_response
        }
    except Exception as error:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PowerPoint generation failed: {error}") from error
